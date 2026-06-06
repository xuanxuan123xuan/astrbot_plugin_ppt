"""
astrbot_plugin_ppt
=================
自动解析 PPTX / PDF / DOCX 课件并生成结构化总结。

触发方式：
  - 向机器人发送课件文件（PPTX/PDF/DOCX），自动提取文本并注入 LLM 请求生成总结
  - /kw   本地 jieba 关键词提取
  - /summarize  手动触发最近一份课件的重新总结
"""

import asyncio
import traceback
from pathlib import Path
from typing import Dict, List, Optional

from astrbot.api.event import filter, AstrMessageEvent
from astrbot.api.star import Context, Star, register
from astrbot.api import logger
from astrbot.api.provider import ProviderRequest

# ---- 文件类型检测依赖的动态导入 ----
# 使用动态导入避免未安装依赖时插件加载失败
# 用户需在 requirements.txt 中声明依赖，AstrBot 会在安装时提示

_MAX_CHARS = 8000  # 注入 LLM 的最大字符数，防止超过上下文限制

# ---- 文本提取函数 ----

def _extract_pptx(file_path: str, max_chars: int = _MAX_CHARS) -> str:
    """从 PPTX 提取文本，按页组织，超出截断"""
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    total = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
                        total += len(t)
        if slide_parts:
            texts.append(f"【第{slide_num}页】" + " | ".join(slide_parts))
        if total >= max_chars:
            texts.append("...（内容过长，后续页面已截断）")
            break
    return "\n".join(texts)


def _extract_pdf(file_path: str, max_chars: int = _MAX_CHARS) -> str:
    """从 PDF 提取文本，按页组织，超出截断"""
    import fitz  # pymupdf

    doc = fitz.open(file_path)
    texts = []
    total = 0
    try:
        for page_num in range(doc.page_count):
            page = doc[page_num]
            page_text = page.get_text("text").strip()
            if page_text:
                texts.append(f"【第{page_num + 1}页】" + page_text)
                total += len(page_text)
            if total >= max_chars:
                texts.append("...（内容过长，后续页面已截断）")
                break
    finally:
        doc.close()
    return "\n".join(texts)


def _extract_docx(file_path: str, max_chars: int = _MAX_CHARS) -> str:
    """从 DOCX 提取文本，按段落组织，超出截断"""
    from docx import Document

    doc = Document(file_path)
    texts = []
    total = 0
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            texts.append(t)
            total += len(t)
        if total >= max_chars:
            texts.append("...（内容过长，后续段落已截断）")
            break
    return "\n".join(texts)


def extract_text(file_path: str, max_chars: int = _MAX_CHARS) -> str:
    """根据扩展名分派到对应的提取函数"""
    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        return _extract_pptx(file_path, max_chars)
    elif ext == ".pdf":
        return _extract_pdf(file_path, max_chars)
    elif ext == ".docx":
        return _extract_docx(file_path, max_chars)
    elif ext == ".ppt":
        # 旧版 PPT 通常也能被 python-pptx 处理
        return _extract_pptx(file_path, max_chars)
    else:
        raise ValueError(f"不支持的文件格式: {ext}（仅支持 pptx / pdf / docx / ppt）")


# ---- 插件主类 ----

@register(
    "astrbot_plugin_ppt",
    "wuxuan",
    "自动解析 PPT/PDF/DOCX 课件并生成结构化总结，支持关键词分析",
    "1.0.0",
    "https://github.com/wuxuan/astrbot_plugin_ppt",
)
class CoursewareSummarizer(Star):
    def __init__(self, context: Context):
        super().__init__(context)
        # 待注入 LLM 请求的文件内容队列
        self._pending: List[Dict] = []
        # 最近一份课件文本（供 /kw 和 /summarize 使用）
        self._last_text: Optional[str] = None
        self._last_name: Optional[str] = None

    # ====================
    #  文件上传检测
    # ====================

    @filter.event_message_type(filter.EventMessageType.ALL)
    async def on_receive_msg(self, event: AstrMessageEvent):
        """监听所有消息，检测 PPT/PDF/DOCX 文件上传"""
        if not event.is_at_or_wake_command:
            return

        for item in event.message_obj.message:
            # 需要导入消息组件检测 File 类型
            from astrbot.api import message_components as Comp

            if not isinstance(item, Comp.File):
                continue

            # 检查扩展名
            name = getattr(item, "name", None) or "unknown"
            ext = Path(name).suffix.lower()
            if ext not in (".pptx", ".pdf", ".docx", ".ppt"):
                logger.info(f"跳过非课件文件: {name}")
                continue

            try:
                file_path = await item.get_file()
                logger.info(f"收到课件文件: {name} -> {file_path}")

                content = extract_text(file_path, _MAX_CHARS)

                self._pending.append({
                    "name": name,
                    "content": content,
                    "length": len(content),
                })
                self._last_text = content
                self._last_name = name

                logger.info(f"课件提取完成: {name}, {len(content)} 字符")

            except ValueError as e:
                logger.warning(f"{e}")
            except ImportError as e:
                logger.error(f"缺少依赖，无法解析课件: {e}")
            except Exception:
                logger.error(f"解析课件失败:\n{traceback.format_exc()}")

    # ====================
    #  LLM 请求注入（自动总结）
    # ====================

    @filter.on_llm_request()
    async def on_request(self, event: AstrMessageEvent, req: ProviderRequest):
        """将课件文本注入 LLM 请求，引导模型生成结构化总结"""
        if not self._pending:
            return

        # 取最早的一份（通常只有一份）
        info = self._pending.pop(0)

        # 清空队列避免重复注入（极端并发场景下安全处理）
        if self._pending:
            logger.warning(f"仍有 {len(self._pending)} 份课件未处理，已丢弃")
        self._pending.clear()

        # 构造课件上下文提示
        summary_prompt = (
            f"\n\n[SYSTEM INSTRUCTION - HIGH PRIORITY]\n"
            f"用户刚刚上传了一份课件文件，请**优先**对其进行结构化总结，"
            f"再处理用户可能附带的其他请求。\n\n"
            f"=== 课件文件信息 ===\n"
            f"文件名: {info['name']}\n"
            f"文本长度: {info['length']} 字符\n\n"
            f"=== 课件正文 ===\n"
            f"{info['content']}\n"
            f"=== 正文结束 ===\n\n"
            f"请在回复中输出以下内容（使用 Markdown 排版）：\n"
            f"## 一、核心主题与学科领域\n"
            f"## 二、关键知识点（3-8 个要点）\n"
            f"## 三、内容结构大纲\n"
            f"## 四、一句话总结\n"
        )

        # 追加到用户输入末尾，让 LLM 作为本轮对话的一部分处理
        req.prompt += summary_prompt

        logger.info(f"已将课件 {info['name']} 注入本轮 LLM 请求")

    # ====================
    #  /kw 本地关键词分析
    # ====================

    @filter.command("kw")
    async def keyword_analysis(self, event: AstrMessageEvent):
        """本地 jieba 关键词提取"""
        if not self._last_text:
            yield event.plain_result(
                "尚未收到课件文件。请先发送 PPT / PDF / DOCX 文件。"
            )
            return

        try:
            import jieba.analyse
        except ImportError:
            yield event.plain_result(
                "缺少 jieba 依赖。请在插件目录下执行: pip install jieba"
            )
            return

        keywords = jieba.analyse.extract_tags(
            self._last_text, topK=15, withWeight=True
        )

        lines = [f"**课件关键词分析**（{self._last_name}）\n"]
        for kw, weight in keywords:
            lines.append(f"- {kw}  `{weight:.3f}`")

        yield event.plain_result("\n".join(lines))

    # ====================
    #  /summarize 手动重总结
    # ====================

    @filter.command("summarize")
    async def manual_summarize(self, event: AstrMessageEvent):
        """手动将最近课件重新送入 LLM 总结"""
        if not self._last_text:
            yield event.plain_result(
                "尚未收到课件文件。请先发送 PPT / PDF / DOCX 文件。"
            )
            return

        # 将最近课件放入待注入队列，本轮 LLM 请求会触发 on_request
        self._pending.append({
            "name": self._last_name or "未知文件",
            "content": self._last_text,
            "length": len(self._last_text),
        })

        yield event.plain_result(
            f"已重新加载课件「{self._last_name}」，正在生成总结..."
        )

    # ====================
    #  清理
    # ====================

    async def terminate(self):
        """插件卸载/停用时清理"""
        self._pending.clear()
        self._last_text = None
        self._last_name = None
        logger.info("课件总结插件已卸载")